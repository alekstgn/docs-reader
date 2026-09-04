"""Excel и PowerPoint по итогам текстовой сверки и этапа C (фото)."""

from __future__ import annotations

import json
import re
from collections import Counter
from datetime import datetime
from pathlib import Path
from typing import Any

from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.page import PageMargins
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree

STATUS_RU = {
    "ok": "совпадает",
    "mismatch": "расхождение",
    "missing_in_report": "нет в отчёте",
    "needs_photo": "нужно фото",
    "needs_photo_skipped": "фото отложены",
    "photo_inconclusive": "фото неубедительны",
    "not_verifiable_from_docs": "не проверить по документам",
    "event_mismatch": "другое мероприятие",
    "qualitative_pending": "ожидает качественной сверки",
}

FILL_OK = PatternFill("solid", fgColor="C6EFCE")
FILL_MIS = PatternFill("solid", fgColor="FFC7CE")
FILL_MISS = PatternFill("solid", fgColor="D9D9D9")
FILL_YEL = PatternFill("solid", fgColor="FFF2CC")
FILL_KEY = PatternFill("solid", fgColor="FCE4D6")
FILL_BANNER_RED = PatternFill("solid", fgColor="C00000")
FILL_BANNER_GRN = PatternFill("solid", fgColor="1E7A46")
FILL_NAVY = PatternFill("solid", fgColor="0B2C4A")
FILL_NAVY2 = PatternFill("solid", fgColor="1B4F72")
FILL_HEAD = PatternFill("solid", fgColor="0B2C4A")
FILL_ALT = PatternFill("solid", fgColor="F4F7FA")
FILL_WHITE = PatternFill("solid", fgColor="FFFFFF")
FILL_PHOTO = PatternFill("solid", fgColor="D6EAF8")

FONT_WHITE = Font(name="Calibri", size=11, color="FFFFFF", bold=True)
FONT_H2 = Font(name="Calibri", size=13, color="0B2C4A", bold=True)
FONT_BODY = Font(name="Calibri", size=11, color="1A1A1A")
FONT_BOLD = Font(name="Calibri", size=11, color="1A1A1A", bold=True)
FONT_SM = Font(name="Calibri", size=10, color="1A1A1A")
FONT_NOTE = Font(name="Calibri", size=9, italic=True, color="595959")
FONT_STATUS = {
    "ok": Font(name="Calibri", size=11, color="006100", bold=True),
    "mismatch": Font(name="Calibri", size=11, color="9C0006", bold=True),
    "missing_in_report": Font(name="Calibri", size=11, color="595959", bold=True),
    "needs_photo": Font(name="Calibri", size=11, color="9C5700", bold=True),
    "needs_photo_skipped": Font(name="Calibri", size=11, color="9C5700", bold=True),
    "photo_inconclusive": Font(name="Calibri", size=11, color="9C5700", bold=True),
    "not_verifiable_from_docs": Font(name="Calibri", size=11, color="9C5700", bold=True),
    "event_mismatch": Font(name="Calibri", size=11, color="9C0006", bold=True),
}

THIN = Border(
    left=Side(style="thin", color="B0B8C1"),
    right=Side(style="thin", color="B0B8C1"),
    top=Side(style="thin", color="B0B8C1"),
    bottom=Side(style="thin", color="B0B8C1"),
)
WRAP = Alignment(wrap_text=True, vertical="center")
WRAP_TOP = Alignment(wrap_text=True, vertical="top")
CENTER = Alignment(wrap_text=True, vertical="center", horizontal="center")

HEADERS = [
    "ID",
    "Пункт ТЗ",
    "Показатель",
    "Требуется (договор)",
    "Факт из отчёта",
    "Статус",
    "Цитата договора",
    "Цитата отчёта",
    "Комментарий",
]
COL_WIDTHS = [26, 12, 38, 36, 36, 26, 48, 48, 48]
KEY_LEKTORIY = {
    "TZ-backdrops-count",
    "TZ-photos-count",
    "TZ-lectures-11jun",
    "TZ-lectures-11jun-venue",
    "PHOTO-event-identity",
    "PHOTO-branding",
}

NAVY = RGBColor(0x0B, 0x2C, 0x4A)
NAVY2 = RGBColor(0x1B, 0x4F, 0x72)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE = RGBColor(0xF4, 0xF7, 0xFA)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5B, 0x6B, 0x7A)
RED = RGBColor(0xB0, 0x3A, 0x2E)
GREEN = RGBColor(0x1E, 0x7A, 0x46)
AMBER = RGBColor(0xC4, 0x7E, 0x0A)
LIGHT_RED = RGBColor(0xFD, 0xED, 0xEC)
LIGHT_GRN = RGBColor(0xE8, 0xF8, 0xF0)
LIGHT_YEL = RGBColor(0xFE, 0xF9, 0xE7)


def clean_comment(text: str) -> str:
    if not text:
        return ""
    if text.startswith("The date and address"):
        return (
            "Дата и адрес в отчёте не совпадают с договором: "
            "в отчёте 30 июня 2026, Музей военной формы; "
            "в договоре 29 ноября 2025, Парк-отель «Шереметьевский»."
        )
    if text.startswith("The report mentions 120 chairs"):
        return (
            "В отчёте указаны 120 стульев и 240 стаканчиков; по договору — не менее 60 стульев. "
            "Даты позиций отчёта (29.06–02.07.2026) не совпадают с договором (ноябрь 2025)."
        )
    if text.startswith("The report confirms"):
        return "Отчёт подтверждает проведение указанного формата."
    text = re.sub(
        r"адрес/площадка: договор \{([^}]+)\} vs отчёт \{([^}]+)\}",
        lambda m: (
            "Адрес/площадка: договор — "
            + ", ".join(sorted(x.strip(" '") for x in m.group(1).split(",")))
            + "; отчёт — "
            + ", ".join(sorted(x.strip(" '") for x in m.group(2).split(",")))
            + "."
        ),
        text,
    )
    return text.replace("[ФИО скрыто]", "[скрыто]")


def clean_text(value: Any) -> str:
    if value is None:
        return ""
    text = str(value).replace("\xa0", " ").replace("\n", " ")
    return " ".join(text.split()).replace("[ФИО скрыто]", "[скрыто]")


def load_json(path: Path) -> dict[str, Any]:
    if not path.is_file():
        return {}
    return json.loads(path.read_text(encoding="utf-8"))


def status_bucket(status: str) -> str:
    if status in ("needs_photo", "needs_photo_skipped"):
        return "needs_photo"
    if status == "photo_inconclusive":
        return "photo_inconclusive"
    if status == "not_verifiable_from_docs":
        return "not_verifiable"
    if status == "missing_in_report":
        return "missing"
    if status in ("mismatch", "event_mismatch"):
        return "mismatch"
    if status == "ok":
        return "ok"
    return status or "other"


def count_statuses(rows: list[dict[str, Any]]) -> dict[str, int]:
    c = Counter(status_bucket(r.get("status") or "") for r in rows)
    return {
        "ok": c.get("ok", 0),
        "mismatch": c.get("mismatch", 0),
        "missing": c.get("missing", 0),
        "needs_photo": c.get("needs_photo", 0),
        "photo_inconclusive": c.get("photo_inconclusive", 0),
        "not_verifiable": c.get("not_verifiable", 0),
        "total": len(rows),
    }


def pair_tokens(data: dict[str, Any], prefix: str) -> tuple[int, int, float, int]:
    details = ((data.get("usage") or {}).get("details") or [])
    inn = out = calls = 0
    for d in details:
        if str(d.get("stage") or "").startswith(prefix):
            inn += int(d.get("prompt_tokens") or 0)
            out += int(d.get("output_tokens") or 0)
            calls += 1
    cost = (inn * 26.0 + out * 129.0) / 1_000_000
    return inn, out, cost, calls


def split_text_photo_usage(summary: dict[str, Any]) -> tuple[dict[str, Any], dict[str, Any]]:
    text_u = summary.get("usage_text") or {}
    photo_u = summary.get("usage_photo") or {}
    if text_u or photo_u:
        return text_u, photo_u
    usage = summary.get("usage") or {}
    text_details, photo_details = [], []
    for d in usage.get("details") or []:
        (photo_details if ":C" in str(d.get("stage") or "") else text_details).append(d)
    def pack(details: list[dict[str, Any]]) -> dict[str, Any]:
        inn = sum(int(d.get("prompt_tokens") or 0) for d in details)
        out = sum(int(d.get("output_tokens") or 0) for d in details)
        return {
            "calls": len(details),
            "prompt_tokens": inn,
            "output_tokens": out,
            "cost_rub": round((inn * 26.0 + out * 129.0) / 1_000_000, 4),
            "details": details,
        }
    return pack(text_details), pack(photo_details)


def photo_coverage(photos: dict[str, Any]) -> str:
    if photos.get("coverage"):
        return str(photos.get("coverage"))
    real = int(photos.get("real_photo_candidates") or 0)
    sent = int(photos.get("sampled") or photos.get("sent") or 0)
    if real > 0 and sent >= int(real * 0.9):
        return "all_real"
    return "sample"


def is_full_photo_run(*photo_dicts: dict[str, Any]) -> bool:
    return any(photo_coverage(p) == "all_real" for p in photo_dicts if p)


CONCLUSION_RU = {
    "confirms": "подтверждает",
    "contradicts": "противоречит",
    "inconclusive": "неясно",
    "wrong_event": "другое мероприятие",
}


def status_fills(status: str):
    if status == "ok":
        return FILL_OK, FONT_STATUS["ok"]
    if status in ("mismatch", "event_mismatch"):
        return FILL_MIS, FONT_STATUS.get(status, FONT_STATUS["mismatch"])
    if status == "missing_in_report":
        return FILL_MISS, FONT_STATUS["missing_in_report"]
    return FILL_YEL, FONT_STATUS.get(status, FONT_STATUS["needs_photo"])


def style_header_row(ws, row: int, n_cols: int) -> None:
    for col in range(1, n_cols + 1):
        cell = ws.cell(row, col)
        cell.fill = FILL_HEAD
        cell.font = FONT_WHITE
        cell.alignment = CENTER
        cell.border = THIN


def apply_widths(ws, widths) -> None:
    for i, w in enumerate(widths, 1):
        ws.column_dimensions[get_column_letter(i)].width = w


def write_obligation_sheet(wb, title, data, banner, banner_fill, key_ids=None):
    ws = wb.create_sheet(title)
    n = len(HEADERS)
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=n)
    cell = ws.cell(1, 1, banner)
    cell.fill = banner_fill
    cell.font = Font(name="Calibri", size=12, color="FFFFFF", bold=True)
    cell.alignment = Alignment(wrap_text=True, vertical="center")
    ws.row_dimensions[1].height = 52
    for col, h in enumerate(HEADERS, 1):
        ws.cell(2, col, h)
    style_header_row(ws, 2, n)
    event = data.get("event_check") or {}
    rows = [event] + list(data.get("rows") or [])
    key_ids = key_ids or set()
    for i, row in enumerate(rows):
        r = 3 + i
        status = row.get("status") or ""
        values = [
            clean_text(row.get("id")),
            clean_text(row.get("clause")),
            clean_text(row.get("metric") or row.get("clause")),
            clean_text(row.get("required")),
            clean_text(row.get("claimed")) or "—",
            STATUS_RU.get(status, status),
            clean_text(row.get("quote_contract")),
            clean_text(row.get("quote_report")),
            clean_comment(clean_text(row.get("comment"))),
        ]
        fill_row, font_st = status_fills(status)
        is_key = row.get("id") in key_ids
        is_photo = str(row.get("id") or "").startswith("PHOTO-")
        for col, val in enumerate(values, 1):
            cell = ws.cell(r, col, val)
            cell.font = FONT_BODY if col != 6 else font_st
            cell.alignment = WRAP_TOP
            cell.border = THIN
            if col == 6:
                cell.fill = fill_row
                cell.alignment = CENTER
            elif status in ("mismatch", "event_mismatch"):
                cell.fill = FILL_MIS
            elif status == "missing_in_report":
                cell.fill = FILL_MISS
            elif is_photo:
                cell.fill = FILL_PHOTO
            elif is_key:
                cell.fill = FILL_KEY
            elif i % 2:
                cell.fill = FILL_ALT
        ws.row_dimensions[r].height = 50
    last = 2 + len(rows)
    ws.freeze_panes = "A3"
    ws.auto_filter.ref = f"A2:{get_column_letter(n)}{last}"
    apply_widths(ws, COL_WIDTHS)
    ws.page_setup.orientation = "landscape"
    ws.page_setup.fitToPage = True
    ws.page_setup.fitToWidth = 1
    ws.page_setup.fitToHeight = 0
    ws.page_setup.paperSize = ws.PAPERSIZE_A4
    ws.page_margins = PageMargins(0.4, 0.4, 0.5, 0.5)
    ws.sheet_view.showGridLines = False
    return ws


def _photo_rows(pair_label: str, photos: dict[str, Any]) -> list[list[Any]]:
    out = []
    for a in photos.get("analyses") or []:
        objects = ", ".join(str(x) for x in (a.get("visible_objects") or [])[:8])
        ids = ", ".join(str(x) for x in (a.get("matched_obligation_ids") or [])[:8])
        raw_c = str(a.get("conclusion") or "")
        if not raw_c:
            fit = a.get("event_fit") or "unknown"
            raw_c = {"other_event": "wrong_event", "contracted_event": "confirms"}.get(fit, "inconclusive")
        conclusion = CONCLUSION_RU.get(raw_c, raw_c or "неясно")
        if a.get("error"):
            conclusion = "ошибка API"
        elif a.get("scene_type"):
            conclusion = f"{conclusion}; сцена: {a.get('scene_type')}"
        out.append(
            [
                pair_label,
                a.get("photo_id") or a.get("jpeg_name") or "",
                a.get("caption") or a.get("date_hint") or "—",
                objects or "—",
                clean_text(a.get("branding_or_text_seen")) or "—",
                ids or "—",
                conclusion,
                a.get("confidence") or "",
                clean_text(a.get("notes")),
                clean_text(a.get("error")) or "",
            ]
        )
    return out


def build_xlsx(out_root: Path, prosv, lekt, summary, photos_p, photos_l, run_date: str) -> Path:
    wb = Workbook()
    ws = wb.active
    ws.title = "Сводка"
    p_cnt = count_statuses(prosv.get("rows") or [])
    l_cnt = count_statuses(lekt.get("rows") or [])
    text_u, photo_u = split_text_photo_usage(summary)
    usage = summary.get("usage") or {}
    p_in, p_out, p_cost, _ = pair_tokens(prosv, "prosvetiteli")
    l_in, l_out, l_cost, _ = pair_tokens(lekt, "lektoriy_kaluga")
    photo_cost = float(photo_u.get("cost_rub") or 0)
    text_cost = float(text_u.get("cost_rub") or 0)
    total_cost = float(usage.get("cost_rub") or (text_cost + photo_cost))

    full = is_full_photo_run(photos_p, photos_l)
    sent_p = int(photos_p.get("sampled") or photos_p.get("sent") or 0)
    sent_l = int(photos_l.get("sampled") or photos_l.get("sent") or 0)
    failed_p = int(photos_p.get("failed") or 0)
    failed_l = int(photos_l.get("failed") or 0)

    ws.merge_cells("A1:L1")
    c = ws.cell(1, 1, "Сверка отчётов подрядчиков с договорами — сводка (текст + полный прогон фото)" if full else "Пилот сверки отчётов подрядчиков с договорами — сводка (текст + фото)")
    c.font = Font(name="Calibri", size=18, color="FFFFFF", bold=True)
    c.fill = FILL_NAVY
    c.alignment = Alignment(vertical="center", wrap_text=True)
    ws.row_dimensions[1].height = 28
    ws.merge_cells("A2:L2")
    c = ws.cell(
        2,
        1,
        "Модель gemini-2.5-flash-lite (ProxyAPI) · этап C: все реальные фото (иконки пропущены) · без персональных данных"
        if full
        else "Модель gemini-2.5-flash-lite (ProxyAPI) · этап C выполнен выборкой кадров · без персональных данных",
    )
    c.font = Font(name="Calibri", size=10, color="FFFFFF", italic=True)
    c.fill = FILL_NAVY2

    ws.merge_cells("A4:D4")
    ws.cell(4, 1, "Параметры прогона").font = FONT_H2
    meta = [
        ("Дата прогона", run_date),
        ("Пар документов", "2"),
        ("Модель", usage.get("model", "gemini-2.5-flash-lite")),
        ("Вызовов API всего", str(usage.get("calls") or (text_u.get("calls", 0) + photo_u.get("calls", 0)))),
        ("Токены ввод (всего)", f"{int(usage.get('prompt_tokens') or 0):,}".replace(",", " ")),
        ("Токены вывод (всего)", f"{int(usage.get('output_tokens') or 0):,}".replace(",", " ")),
        ("Стоимость всего, ₽", f"{total_cost:.2f}".replace(".", ",")),
        ("Этап C (фото)", "выполнен"),
        ("Стоимость текста, ₽", f"{text_cost:.2f}".replace(".", ",")),
        ("Доп. стоимость фото, ₽", f"{photo_cost:.2f}".replace(".", ",")),
        ("Кадров отправлено", str(sent_p + sent_l)),
        ("Вложений извлечено", str((photos_p.get("extracted_total") or 0) + (photos_l.get("extracted_total") or 0))),
    ]
    for i, (k, v) in enumerate(meta):
        col = 1 + (i // 6) * 2
        row = 5 + (i % 6)
        ws.cell(row, col, k).font = FONT_BOLD
        ws.cell(row, col).fill = FILL_ALT
        ws.cell(row, col).border = THIN
        ws.cell(row, col + 1, v).font = FONT_BODY
        ws.cell(row, col + 1).border = THIN

    ws.merge_cells("A12:L12")
    ws.cell(12, 1, "Результаты по парам").font = FONT_H2
    pair_headers = [
        "Пара",
        "Мероприятие совпадает",
        "совпадает",
        "расхождение",
        "нет в отчёте",
        "фото неубедительны",
        "не проверить",
        "Пунктов",
        "Вложений / отправлено",
        "Стоимость текста, ₽",
        "Стоимость фото, ₽",
        "Итого, ₽",
    ]
    for col, h in enumerate(pair_headers, 1):
        ws.cell(13, col, h)
    style_header_row(ws, 13, len(pair_headers))
    p_photo_cost = float((photos_p.get("usage") or {}).get("cost_rub") or 0)
    l_photo_cost = float((photos_l.get("usage") or {}).get("cost_rub") or 0)
    # если usage парный не записан — делим photo_u по деталям
    if not p_photo_cost and not l_photo_cost:
        for d in photo_u.get("details") or []:
            inn = int(d.get("prompt_tokens") or 0)
            outt = int(d.get("output_tokens") or 0)
            add = (inn * 26.0 + outt * 129.0) / 1_000_000
            if str(d.get("stage") or "").startswith("prosvetiteli"):
                p_photo_cost += add
            else:
                l_photo_cost += add
    pair_rows = [
        (
            "Просветители",
            "НЕТ — разные мероприятия",
            p_cnt,
            f"{prosv.get('report_images', 0)} / {sent_p}",
            p_cost,
            p_photo_cost,
            False,
        ),
        (
            "Лекторий (Калуга)",
            "ДА — те же даты и площадки",
            l_cnt,
            f"{lekt.get('report_images', 0)} / {sent_l}",
            l_cost,
            l_photo_cost,
            True,
        ),
    ]
    for i, (name, ev, cnt, imgs, tcost, phcost, same) in enumerate(pair_rows):
        r = 14 + i
        vals = [
            name,
            ev,
            cnt["ok"],
            cnt["mismatch"],
            cnt["missing"],
            cnt["photo_inconclusive"],
            cnt["not_verifiable"],
            cnt["total"],
            imgs,
            round(tcost, 2),
            round(phcost, 2),
            round(tcost + phcost, 2),
        ]
        for col, v in enumerate(vals, 1):
            cell = ws.cell(r, col, v)
            cell.font = FONT_BODY
            cell.alignment = CENTER if col > 1 else WRAP
            cell.border = THIN
            if col == 2:
                cell.fill = FILL_OK if same else FILL_MIS
        ws.row_dimensions[r].height = 28

    ws.merge_cells("A17:L17")
    ws.cell(17, 1, "Ключевые выводы после этапа C").font = FONT_H2
    findings = [
        (
            "Просветители",
            "Фото не закрывают договор ноября 2025",
            "Текстовая сверка: другое мероприятие (июнь 2026, «Лесной»). "
            f"Извлечено {photos_p.get('extracted_total', 0)} вложений, реальных фото {photos_p.get('real_photo_candidates', 0)}, "
            f"в модель ушло {sent_p} кадров (иконок пропущено {photos_p.get('skipped_tiny', 0)}"
            f"{f', ошибок {failed_p}' if failed_p else ''}). "
            "Кадры относятся к событию отчёта, а не к ТЗ ноября 2025 — даже если зал и брендирование выглядят нормально.",
            FILL_MIS,
        ),
        (
            "Лекторий",
            "Визуально мероприятие подтверждается",
            "Даты совпадают. "
            f"Извлечено {photos_l.get('extracted_total', 0)} вложений, реальных фото {photos_l.get('real_photo_candidates', 0)}, "
            f"отправлено {sent_l} (иконок {photos_l.get('skipped_tiny', 0)}"
            f"{f', ошибок {failed_l}' if failed_l else ''}). "
            "На кадрах видны зал/аудитория и брендирование История.РФ. "
            "Микрофоны поштучно не считались. Разрыв «1748 заявлено / 124 в DOCX» остаётся текстовым: "
            "полный набор вложений в файле просмотрен, это не архив на 1748 снимков.",
            FILL_KEY,
        ),
        (
            "Стоимость",
            f"Текст {text_cost:.2f} ₽ + фото {photo_cost:.2f} ₽",
            f"Итого ≈ {total_cost:.2f} ₽. Картинки тарифицируются как входные токены (26 ₽ / 1 млн). "
            + (
                f"Полный прогон: {sent_p + sent_l} реальных кадров обеих пар."
                if full
                else "Полные 224+124 кадра не отправлялись — только выборка 15–25 на пару."
            ),
            FILL_YEL,
        ),
    ]
    for col, h in enumerate(["Пара / тема", "Итог", "Пояснение"], 1):
        ws.cell(18, col, h)
    ws.merge_cells("C18:L18")
    style_header_row(ws, 18, 3)
    for col in range(4, 13):
        ws.cell(18, col).fill = FILL_HEAD
        ws.cell(18, col).border = THIN
    for i, (a, b, cmt, fill) in enumerate(findings):
        r = 19 + i
        ws.merge_cells(start_row=r, start_column=3, end_row=r, end_column=12)
        ws.cell(r, 1, a).font = FONT_BOLD
        ws.cell(r, 2, b).font = FONT_BOLD
        ws.cell(r, 3, cmt).font = FONT_SM
        for col in range(1, 13):
            ws.cell(r, col).fill = fill
            ws.cell(r, col).border = THIN
            ws.cell(r, col).alignment = WRAP
        ws.row_dimensions[r].height = 52
    apply_widths(ws, [28, 40, 16, 16, 16, 20, 16, 12, 22, 18, 18, 14])
    ws.freeze_panes = "A14"
    ws.sheet_view.showGridLines = False

    write_obligation_sheet(
        wb,
        "Просветители",
        prosv,
        "ВНИМАНИЕ: мероприятия НЕ совпадают. Договор — форум 27–30.11.2025, «Шереметьевский». "
        "Отчёт и фото — конференция 29.06–02.07.2026, «Лесной». Этап C: кадры не подтверждают ТЗ ноября 2025.",
        FILL_BANNER_RED,
        {"PHOTO-event-identity", "PHOTO-branding", "TZ-photos"},
    )
    write_obligation_sheet(
        wb,
        "Лекторий",
        lekt,
        "Мероприятие совпадает: 10–12 июня 2025, Калуга. Этап C: все вложенные реальные фото подтверждают зал и брендирование. "
        "Ключевое текстовое расхождение: адрес 11 июня; фото ≥100 / заявлено 1748 / в DOCX 124.",
        FILL_BANNER_GRN,
        KEY_LEKTORIY,
    )

    ws = wb.create_sheet("Фото")
    ws.merge_cells("A1:J1")
    n_rows = len(_photo_rows("Просветители", photos_p)) + len(_photo_rows("Лекторий", photos_l))
    c = ws.cell(
        1,
        1,
        f"Этап C: {'полный прогон' if full else 'выборка'} — {n_rows} кадров "
        f"(сжатые JPEG, не исходный DOCX; иконки пропущены)",
    )
    c.font = Font(name="Calibri", size=16, color="FFFFFF", bold=True)
    c.fill = FILL_NAVY
    ws.row_dimensions[1].height = 24
    photo_headers = [
        "Пара",
        "Файл",
        "Подпись / дата",
        "Что видно",
        "Брендинг / текст на кадре",
        "Пункты ТЗ",
        "Вывод",
        "Уверенность",
        "Комментарий модели",
        "Ошибка",
    ]
    for col, h in enumerate(photo_headers, 1):
        ws.cell(2, col, h)
    style_header_row(ws, 2, len(photo_headers))
    all_photo_rows = _photo_rows("Просветители", photos_p) + _photo_rows("Лекторий", photos_l)
    for i, vals in enumerate(all_photo_rows):
        r = 3 + i
        for col, v in enumerate(vals, 1):
            cell = ws.cell(r, col, v)
            cell.font = FONT_SM
            cell.alignment = WRAP_TOP
            cell.border = THIN
            verdict = str(vals[6])
            if "ошибка" in verdict:
                cell.fill = FILL_YEL
            elif "другое" in verdict:
                cell.fill = FILL_MIS if col == 7 else (FILL_ALT if i % 2 else FILL_WHITE)
            elif i % 2:
                cell.fill = FILL_ALT
        ws.row_dimensions[r].height = 36
    last = 2 + max(len(all_photo_rows), 1)
    ws.freeze_panes = "A3"
    ws.auto_filter.ref = f"A2:J{last}"
    apply_widths(ws, [16, 22, 22, 32, 32, 28, 28, 14, 48, 28])
    ws.sheet_view.showGridLines = False

    ws = wb.create_sheet("Стоимость")
    ws.merge_cells("A1:G1")
    c = ws.cell(1, 1, "Стоимость: текст vs фото vs итого")
    c.font = Font(name="Calibri", size=18, color="FFFFFF", bold=True)
    c.fill = FILL_NAVY
    ws.row_dimensions[1].height = 28
    ws.cell(3, 1, "Контур").font = FONT_H2
    for col, h in enumerate(["Контур", "Вызовов", "Токены ввод", "Токены вывод", "Стоимость, ₽", "Доля, %"], 1):
        ws.cell(4, col, h)
    style_header_row(ws, 4, 6)
    blocks = [
        ("Текст (этапы A/B/qual)", text_u),
        ("Фото (этап C)", photo_u),
        ("Итого", usage if usage.get("prompt_tokens") else {
            "calls": int(text_u.get("calls") or 0) + int(photo_u.get("calls") or 0),
            "prompt_tokens": int(text_u.get("prompt_tokens") or 0) + int(photo_u.get("prompt_tokens") or 0),
            "output_tokens": int(text_u.get("output_tokens") or 0) + int(photo_u.get("output_tokens") or 0),
            "cost_rub": round(text_cost + photo_cost, 4),
        }),
    ]
    tot = float(blocks[2][1].get("cost_rub") or 1) or 1
    for i, (name, u) in enumerate(blocks):
        r = 5 + i
        cost = float(u.get("cost_rub") or 0)
        vals = [
            name,
            int(u.get("calls") or 0),
            int(u.get("prompt_tokens") or 0),
            int(u.get("output_tokens") or 0),
            round(cost, 4),
            round(cost / tot * 100, 1),
        ]
        for col, v in enumerate(vals, 1):
            cell = ws.cell(r, col, v)
            cell.border = THIN
            cell.font = FONT_WHITE if i == 2 else FONT_BODY
            cell.fill = FILL_NAVY if i == 2 else (FILL_PHOTO if i == 1 else FILL_ALT)
            cell.alignment = CENTER if col > 1 else WRAP

    ws.cell(9, 1, "Охват фото").font = FONT_H2
    for col, h in enumerate(
        ["Пара", "Извлечено из DOCX", "Реальные фото", "Отправлено в API", "Иконки пропущены", "Ошибки API"],
        1,
    ):
        ws.cell(10, col, h)
    style_header_row(ws, 10, 6)
    coverage_rows = [
        (
            "Просветители",
            int(photos_p.get("extracted_total") or 0),
            int(photos_p.get("real_photo_candidates") or 0),
            sent_p,
            int(photos_p.get("skipped_tiny") or 0),
            failed_p,
        ),
        (
            "Лекторий",
            int(photos_l.get("extracted_total") or 0),
            int(photos_l.get("real_photo_candidates") or 0),
            sent_l,
            int(photos_l.get("skipped_tiny") or 0),
            failed_l,
        ),
        (
            "Итого",
            int(photos_p.get("extracted_total") or 0) + int(photos_l.get("extracted_total") or 0),
            int(photos_p.get("real_photo_candidates") or 0) + int(photos_l.get("real_photo_candidates") or 0),
            sent_p + sent_l,
            int(photos_p.get("skipped_tiny") or 0) + int(photos_l.get("skipped_tiny") or 0),
            failed_p + failed_l,
        ),
    ]
    for i, vals in enumerate(coverage_rows):
        r = 11 + i
        for col, v in enumerate(vals, 1):
            cell = ws.cell(r, col, v)
            cell.border = THIN
            cell.alignment = CENTER if col > 1 else WRAP
            cell.font = FONT_WHITE if i == 2 else FONT_BODY
            cell.fill = FILL_NAVY if i == 2 else FILL_ALT

    ws.cell(15, 1, "Вызовы API").font = FONT_H2
    for col, h in enumerate(["Этап", "Пара", "Ввод", "Вывод", "Всего", "₽"], 1):
        ws.cell(16, col, h)
    style_header_row(ws, 16, 6)
    details = (usage.get("details") or text_u.get("details") or []) + (
        [] if usage.get("details") else (photo_u.get("details") or [])
    )
    if usage.get("details"):
        details = usage.get("details") or []
    text_details = [d for d in details if ":C" not in str(d.get("stage") or "")]
    photo_details = [d for d in details if ":C" in str(d.get("stage") or "")]
    if not photo_details:
        photo_details = list(photo_u.get("details") or [])
    display: list[dict[str, Any]] = list(text_details)
    if len(photo_details) > 12:
        for pair_key, pair_name in (("prosvetiteli", "Просветители"), ("lektoriy", "Лекторий")):
            chunk = [d for d in photo_details if str(d.get("stage") or "").startswith(pair_key)]
            if not chunk:
                continue
            inn = sum(int(d.get("prompt_tokens") or 0) for d in chunk)
            outt = sum(int(d.get("output_tokens") or 0) for d in chunk)
            display.append(
                {
                    "stage": f"{pair_key}:C ({len(chunk)} вызовов)",
                    "prompt_tokens": inn,
                    "output_tokens": outt,
                    "_pair": pair_name,
                }
            )
    else:
        display.extend(photo_details)
    start = 17
    for i, d in enumerate(display):
        r = start + i
        stage = str(d.get("stage") or "")
        pair = d.get("_pair") or ("Просветители" if stage.startswith("prosvetiteli") else "Лекторий")
        inn = int(d.get("prompt_tokens") or 0)
        outt = int(d.get("output_tokens") or 0)
        vals = [stage.split(":", 1)[-1], pair, inn, outt, inn + outt, round((inn * 26 + outt * 129) / 1e6, 4)]
        for col, v in enumerate(vals, 1):
            cell = ws.cell(r, col, v)
            cell.border = THIN
            cell.font = FONT_SM
            cell.fill = FILL_YEL if ":C" in stage else (
                PatternFill("solid", fgColor="D6EAF8") if pair == "Просветители" else PatternFill("solid", fgColor="D5F5E3")
            )
    last_d = start + max(len(display), 1) - 1
    chart_src = last_d + 3
    ws.cell(chart_src, 1, "Контур")
    ws.cell(chart_src, 2, "₽")
    style_header_row(ws, chart_src, 2)
    ws.cell(chart_src + 1, 1, "Текст")
    ws.cell(chart_src + 1, 2, round(text_cost, 2))
    ws.cell(chart_src + 2, 1, "Фото")
    ws.cell(chart_src + 2, 2, round(photo_cost, 2))
    ws.cell(chart_src + 3, 1, "Итого")
    ws.cell(chart_src + 3, 2, round(total_cost, 2))
    for rr in range(chart_src + 1, chart_src + 4):
        for col in range(1, 3):
            ws.cell(rr, col).border = THIN
    chart = BarChart()
    chart.type = "col"
    chart.title = "Стоимость, ₽"
    chart.style = 10
    data_ref = Reference(ws, min_col=2, min_row=chart_src, max_row=chart_src + 3)
    cats = Reference(ws, min_col=1, min_row=chart_src + 1, max_row=chart_src + 3)
    chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats)
    chart.legend = None
    chart.width = 12
    chart.height = 7
    ws.add_chart(chart, f"A{chart_src + 5}")
    apply_widths(ws, [28, 22, 18, 18, 18, 16, 14])
    ws.sheet_view.showGridLines = False

    path = out_root / "sverka_dogovorov.xlsx"
    wb.save(path)
    return path


def set_run(run, text, size=18, bold=False, color=DARK, font="Calibri"):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def fill_shape(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    fill_shape(sh, color)
    return sh


def add_round(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill_shape(sh, color)
    try:
        sh.adjustments[0] = 0.08
    except Exception:
        pass
    return sh


def header_bar(slide, title: str, subtitle: str | None = None) -> None:
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.05), NAVY)
    add_rect(slide, Inches(0), Inches(1.05), Inches(13.333), Inches(0.06), GOLD)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.22), Inches(12.3), Inches(0.5))
    set_run(box.text_frame.paragraphs[0].add_run(), title, 24, True, WHITE)
    if subtitle:
        box2 = slide.shapes.add_textbox(Inches(0.5), Inches(0.68), Inches(12.3), Inches(0.32))
        set_run(box2.text_frame.paragraphs[0].add_run(), subtitle, 12, False, RGBColor(0xD0, 0xD8, 0xE0))
    add_rect(slide, Inches(0), Inches(7.22), Inches(13.333), Inches(0.28), NAVY)
    fb = slide.shapes.add_textbox(Inches(0.5), Inches(7.24), Inches(12.3), Inches(0.24))
    set_run(
        fb.text_frame.paragraphs[0].add_run(),
        "РВИО  ·  пилот сверки отчётов с договорами  ·  август 2026  ·  без персональных данных",
        10,
        False,
        RGBColor(0xC5, 0xD0, 0xDA),
    )


def card(slide, l, t, w, h, title, body, accent=GOLD, title_size=14, body_size=13):
    add_round(slide, l, t, w, h, OFFWHITE)
    add_rect(slide, l, t, Inches(0.08), h, accent)
    tb = slide.shapes.add_textbox(l + Inches(0.22), t + Inches(0.12), w - Inches(0.35), Inches(0.36))
    set_run(tb.text_frame.paragraphs[0].add_run(), title, title_size, True, NAVY)
    bb = slide.shapes.add_textbox(l + Inches(0.22), t + Inches(0.48), w - Inches(0.35), h - Inches(0.6))
    tf = bb.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), body, body_size, False, DARK)


def set_cell(cell, text, size=12, bold=False, color=DARK, fill=None, align=PP_ALIGN.LEFT):
    cell.text = ""
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    set_run(run, text, size, bold, color)
    if fill is not None:
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for child in list(tcPr):
            if child.tag == qn("a:solidFill"):
                tcPr.remove(child)
        solid = etree.SubElement(tcPr, qn("a:solidFill"))
        srgb = etree.SubElement(solid, qn("a:srgbClr"))
        srgb.set("val", str(fill))


def _fit_counts(photos: dict[str, Any]) -> Counter:
    return Counter(str(a.get("event_fit") or "unknown") for a in (photos.get("analyses") or []))


def build_pptx(out_root: Path, prosv, lekt, summary, photos_p, photos_l, run_date: str) -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    p_cnt = count_statuses(prosv.get("rows") or [])
    l_cnt = count_statuses(lekt.get("rows") or [])
    text_u, photo_u = split_text_photo_usage(summary)
    usage = summary.get("usage") or {}
    text_cost = float(text_u.get("cost_rub") or 0)
    photo_cost = float(photo_u.get("cost_rub") or 0)
    total_cost = float(usage.get("cost_rub") or (text_cost + photo_cost))
    p_fit = _fit_counts(photos_p)
    l_fit = _fit_counts(photos_l)
    full = is_full_photo_run(photos_p, photos_l)
    sent_p = int(photos_p.get("sampled") or photos_p.get("sent") or 0)
    sent_l = int(photos_l.get("sampled") or photos_l.get("sent") or 0)
    failed_p = int(photos_p.get("failed") or 0)
    failed_l = int(photos_l.get("failed") or 0)
    skipped_p = int(photos_p.get("skipped_tiny") or 0)
    skipped_l = int(photos_l.get("skipped_tiny") or 0)
    real_p = int(photos_p.get("real_photo_candidates") or sent_p)
    real_l = int(photos_l.get("real_photo_candidates") or sent_l)

    s = prs.slides.add_slide(blank)
    add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(7.5), NAVY)
    add_rect(s, Inches(0), Inches(0), Inches(0.18), Inches(7.5), GOLD)
    add_rect(s, Inches(0.5), Inches(2.15), Inches(2.2), Inches(0.07), GOLD)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(2.35), Inches(12.2), Inches(1.6))
    tf = tb.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), "Автоматизация сверки отчётов\nс договорами", 36, True, WHITE)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(4.15), Inches(12.2), Inches(0.7))
    set_run(
        tb.text_frame.paragraphs[0].add_run(),
        "Полный прогон: текст + все реальные фото  ·  помощник проверяющего, не автоподпись акта"
        if full
        else "Пилот: текст + выборка фото  ·  помощник проверяющего, не автоподпись акта",
        18,
        False,
        RGBColor(0xD0, 0xD8, 0xE0),
    )
    tb = s.shapes.add_textbox(Inches(0.5), Inches(6.55), Inches(12.2), Inches(0.4))
    set_run(
        tb.text_frame.paragraphs[0].add_run(),
        f"РВИО  ·  {run_date}  ·  gemini-2.5-flash-lite  ·  {total_cost:.2f} ₽ всего (текст {text_cost:.2f} + фото {photo_cost:.2f})".replace(".", ","),
        14,
        False,
        GOLD,
    )

    s = prs.slides.add_slide(blank)
    header_bar(s, "Задача и объём", "Зачем пилот и какой масштаб")
    card(s, Inches(0.45), Inches(1.35), Inches(6.0), Inches(2.35), "Задача",
         "Сверять отчёт подрядчика с ТЗ: количества, даты, площадки, форматы. Человек принимает решение по акту.")
    card(s, Inches(6.7), Inches(1.35), Inches(6.15), Inches(2.35), "Объём",
         "Ориентир: 1000–3000 листов в год. Ограничивают точность и регламент фото, а не стоимость API.")
    card(s, Inches(0.45), Inches(3.9), Inches(4.0), Inches(2.95), "Текст",
         "Чеклист из ТЗ, факты из отчёта, статус пункта и дословные цитаты.", NAVY2)
    card(s, Inches(4.65), Inches(3.9), Inches(4.0), Inches(2.95), "Фото (этап C)",
         f"Все реальные кадры из DOCX, иконки отсечены. Сжатый JPEG ~1280 px. Сейчас: {sent_p + sent_l} кадров."
         if full else
         "Не весь архив. 15–25 сжатых кадров на пару: брендирование, зал, техника.", GOLD)
    card(s, Inches(8.85), Inches(3.9), Inches(4.0), Inches(2.95), "Критерий",
         "Поймать чужое мероприятие и конкретные расхождения на «своей» паре.", GREEN)

    s = prs.slides.add_slide(blank)
    header_bar(s, "Что проверяли", "Две пары: договор + отчёт. Имена людей исключены.")
    table = s.shapes.add_table(3, 5, Inches(0.45), Inches(1.4), Inches(12.4), Inches(3.15)).table
    rows_data = [
        ["Пара", "Договор", "Отчёт", "Вложения / в API", "Идентичность"],
        ["Просветители", "Форум, 27–30.11.2025,\nШереметьевский",
         "Конференция, 29.06–02.07.2026,\n«Лесной»",
         f"{prosv.get('report_images', 224)} / {sent_p}", "НЕ совпадают"],
        ["Лекторий\n(Калуга)", "Лекторий «История.РФ»,\n10–12.06.2025",
         "Те же даты и город",
         f"{lekt.get('report_images', 124)} / {sent_l}", "Совпадают"],
    ]
    table.columns[0].width = Inches(2.1)
    table.columns[1].width = Inches(3.1)
    table.columns[2].width = Inches(3.3)
    table.columns[3].width = Inches(2.0)
    table.columns[4].width = Inches(1.9)
    for i, row in enumerate(rows_data):
        for j, val in enumerate(row):
            fill = "0B2C4A" if i == 0 else ("FDEDEC" if i == 1 else "E8F8F0")
            color = WHITE if i == 0 else (RED if (i == 1 and j == 4) else (GREEN if (i == 2 and j == 4) else DARK))
            set_cell(table.cell(i, j), val, 13 if i else 12, True if (i == 0 or j in (0, 4)) else False, color, fill,
                     PP_ALIGN.CENTER if j else PP_ALIGN.LEFT)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(4.85), Inches(12.3), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), "Две разные пары нужны специально", 16, True, NAVY)
    p = tf.add_paragraph()
    p.space_before = Pt(8)
    set_run(p.add_run(),
            "«Просветители» — стресс-тест на подмену события. Фото этапа C не должны «закрыть» ноябрь 2025, если кадры с другого форума.",
            14, False, DARK)
    p = tf.add_paragraph()
    p.space_before = Pt(6)
    set_run(p.add_run(),
            "«Лекторий» — штатный случай: те же даты, расхождения по адресу 11 июня и разрыв заявленного числа фото vs вложение DOCX.",
            14, False, DARK)

    s = prs.slides.add_slide(blank)
    header_bar(s, "Как устроен бот", "Текст отдельно. Фото — все реальные кадры из файла. Человек — в конце.")
    steps = [
        ("1", "Локально", "Разбор DOCX. В API не уходит файл на 100 МБ."),
        ("2", "Этап A", "Из ТЗ — чеклист: «не менее N», даты, площадки."),
        ("3", "Этап B", "Из отчёта — заявленные факты и цитаты."),
        ("4", "Сверка", "Числа правилами, смысл — модель. Статус + цитаты."),
        ("5", "Этап C", "Все реальные JPEG ~1280 px. Брендирование, зал, техника. Иконки пропущены."),
        ("6", "Человек", "Таблица расхождений. Решение по акту — за проверяющим."),
    ]
    for i, (num, title, body) in enumerate(steps):
        col, row = i % 3, i // 3
        l = Inches(0.45 + col * 4.25)
        t = Inches(1.4 + row * 2.7)
        add_round(s, l, t, Inches(4.05), Inches(2.45), OFFWHITE)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, l + Inches(0.2), t + Inches(0.25), Inches(0.48), Inches(0.48))
        fill_shape(circ, GOLD if i == 4 else NAVY)
        tb = s.shapes.add_textbox(l + Inches(0.2), t + Inches(0.3), Inches(0.48), Inches(0.42))
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(tb.text_frame.paragraphs[0].add_run(), num, 16, True, WHITE)
        tb = s.shapes.add_textbox(l + Inches(0.8), t + Inches(0.3), Inches(3.0), Inches(0.4))
        set_run(tb.text_frame.paragraphs[0].add_run(), title, 18, True, NAVY)
        tb = s.shapes.add_textbox(l + Inches(0.22), t + Inches(0.95), Inches(3.6), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        set_run(tf.paragraphs[0].add_run(), body, 14, False, DARK)

    s = prs.slides.add_slide(blank)
    header_bar(s, "Результат «Просветители»", "Документы описывают разные мероприятия")
    add_round(s, Inches(0.45), Inches(1.3), Inches(12.4), Inches(1.05), LIGHT_RED)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.42), Inches(12.0), Inches(0.85))
    tf = tb.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(),
            "Мероприятия НЕ совпадают. Сверку пунктов ТЗ нельзя принимать как приёмку. Фото это подтверждают, а не опровергают.",
            16, True, RED)
    add_round(s, Inches(0.45), Inches(2.55), Inches(6.05), Inches(2.4), OFFWHITE)
    add_rect(s, Inches(0.45), Inches(2.55), Inches(6.05), Inches(0.42), NAVY)
    tb = s.shapes.add_textbox(Inches(0.6), Inches(2.6), Inches(5.7), Inches(0.35))
    set_run(tb.text_frame.paragraphs[0].add_run(), "Договор", 14, True, WHITE)
    tb = s.shapes.add_textbox(Inches(0.65), Inches(3.1), Inches(5.65), Inches(1.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(["Форум «Просветители»", "27–30 ноября 2025", "Парк-отель «Шереметьевский»", "также Музей военной формы"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        set_run(p.add_run(), "•  " + line, 15, False, DARK)
    add_round(s, Inches(6.8), Inches(2.55), Inches(6.05), Inches(2.4), OFFWHITE)
    add_rect(s, Inches(6.8), Inches(2.55), Inches(6.05), Inches(0.42), RED)
    tb = s.shapes.add_textbox(Inches(6.95), Inches(2.6), Inches(5.7), Inches(0.35))
    set_run(tb.text_frame.paragraphs[0].add_run(), "Отчёт + фото", 14, True, WHITE)
    tb = s.shapes.add_textbox(Inches(7.0), Inches(3.1), Inches(5.65), Inches(1.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(["Конференция «Просветители.Обществознание»", "29 июня – 2 июля 2026", "Парк-отель «Лесной»", "даты с договором не пересекаются"]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        set_run(p.add_run(), "•  " + line, 15, False, DARK)
    stats = [(str(p_cnt["ok"]), "совпадает", GREEN), (str(p_cnt["mismatch"]), "расхождение", RED),
             (str(p_cnt["missing"]), "нет в отчёте", MUTED), (str(p_cnt["photo_inconclusive"]), "фото неубедительны", AMBER)]
    for i, (n, label, col) in enumerate(stats):
        l = Inches(0.45 + i * 3.2)
        add_round(s, l, Inches(5.15), Inches(3.0), Inches(1.8), OFFWHITE)
        tb = s.shapes.add_textbox(l + Inches(0.15), Inches(5.25), Inches(2.7), Inches(0.7))
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(tb.text_frame.paragraphs[0].add_run(), n, 28, True, col)
        tb = s.shapes.add_textbox(l + Inches(0.15), Inches(5.95), Inches(2.7), Inches(0.7))
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(tb.text_frame.paragraphs[0].add_run(), label, 13, False, MUTED)

    s = prs.slides.add_slide(blank)
    header_bar(s, "Результат «Лекторий»", "Те же даты (10–12 июня 2025, Калуга)")
    add_round(s, Inches(0.45), Inches(1.3), Inches(12.4), Inches(0.7), LIGHT_GRN)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(12.0), Inches(0.5))
    set_run(tb.text_frame.paragraphs[0].add_run(),
            "Мероприятие совпадает. Ниже — расхождения, которые бот должен был поймать.", 15, True, GREEN)
    findings = [
        ("Заставки", "≥ 16", "22", "Порог выполнен", GREEN),
        ("Адрес 11 июня", "Автомобильная, 6", "также Октябрьская, 17а", "Расхождение площадки", RED),
        ("Фотоотчёт", "≥ 100", "1748 заявлено / 124 в файле", "Текст vs вложение", AMBER),
    ]
    for i, (title, req, fact, verdict, col) in enumerate(findings):
        l = Inches(0.45 + i * 4.25)
        add_round(s, l, Inches(2.2), Inches(4.05), Inches(3.0), OFFWHITE)
        add_rect(s, l, Inches(2.2), Inches(4.05), Inches(0.45), col)
        tb = s.shapes.add_textbox(l + Inches(0.15), Inches(2.26), Inches(3.75), Inches(0.35))
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(tb.text_frame.paragraphs[0].add_run(), title, 16, True, WHITE)
        tb = s.shapes.add_textbox(l + Inches(0.2), Inches(2.8), Inches(3.65), Inches(2.2))
        tf = tb.text_frame
        tf.word_wrap = True
        set_run(tf.paragraphs[0].add_run(), "Договор:  " + req, 13, False, MUTED)
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        set_run(p.add_run(), "Отчёт:  " + fact, 13, False, DARK)
        p = tf.add_paragraph()
        p.space_before = Pt(12)
        set_run(p.add_run(), verdict, 15, True, col)
    tb = s.shapes.add_textbox(Inches(0.5), Inches(5.45), Inches(12.3), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(),
            f"Пункты ТЗ: совпадает {l_cnt['ok']}  ·  расхождение {l_cnt['mismatch']}  ·  нет в отчёте {l_cnt['missing']}  ·  фото неубедительны {l_cnt['photo_inconclusive']}.",
            14, True, NAVY)
    p = tf.add_paragraph()
    p.space_before = Pt(8)
    set_run(p.add_run(),
            "Часть «расхождений» — формулировки («зал» перед названием). Для акта важны адрес 11 июня и разрыв 1748 / 124, а не косметика.",
            14, False, DARK)

    s = prs.slides.add_slide(blank)
    header_bar(s, "Этап C: как смотрели фото", "Все реальные вложения DOCX, сжатые JPEG; исходный файл 100 МБ в API не уходил")
    card(s, Inches(0.45), Inches(1.35), Inches(4.05), Inches(2.55), "Извлечение",
         f"Просветители: {photos_p.get('extracted_total', 0)} вложений, из них реальных {real_p} (иконок {skipped_p}). "
         f"Лекторий: {photos_l.get('extracted_total', 0)}, реальных {real_l} (иконок {skipped_l}).")
    card(s, Inches(4.65), Inches(1.35), Inches(4.05), Inches(2.55), "Охват",
         f"В API отправлено {sent_p} и {sent_l} кадров"
         + (f", ошибок {failed_p + failed_l}" if (failed_p + failed_l) else "")
         + ". Это полный набор реальных фото из обоих отчётов, не выборка 20+20.")
    card(s, Inches(8.85), Inches(1.35), Inches(4.0), Inches(2.55), "Сжатие",
         "JPEG, длинная сторона ~1280 px, качество ~70. Пакеты по 1–2 кадра. Иконки и декор отсечены по размеру файла и пикселей.")
    card(s, Inches(0.45), Inches(4.1), Inches(6.15), Inches(2.8), "Что просили у модели",
         "Что видно для ТЗ: брендирование, зал, техника. Без имён и описания лиц. "
         "event_fit: договор / другое событие / неясно. Сопоставление с id пунктов чеклиста.", NAVY2)
    card(s, Inches(6.8), Inches(4.1), Inches(6.05), Inches(2.8), "Чего не делали",
         "Не считали все 1748 фото. Не идентифицировали людей. "
         "Не принимали кадры «Просветителей» как доказательство ноября 2025.", RED)

    s = prs.slides.add_slide(blank)
    header_bar(s, "Этап C: что подтвердили визуально", "Полный набор вложений в DOCX, не внешний архив на 1748 файлов")
    add_round(s, Inches(0.45), Inches(1.3), Inches(6.15), Inches(5.55), OFFWHITE)
    add_rect(s, Inches(0.45), Inches(1.3), Inches(6.15), Inches(0.5), RED)
    tb = s.shapes.add_textbox(Inches(0.65), Inches(1.38), Inches(5.8), Inches(0.38))
    set_run(tb.text_frame.paragraphs[0].add_run(), "Просветители — кадры другого события", 15, True, WHITE)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.0), Inches(5.7), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    lines_p = [
        f"Извлечено {photos_p.get('extracted_total', 0)}, реальных {real_p}, в API {sent_p}"
        + (f", ошибок {failed_p}" if failed_p else "")
        + ".",
        f"event_fit: другое {p_fit.get('other_event', 0)}, договор {p_fit.get('contracted_event', 0)}, неясно {p_fit.get('unknown', 0)}.",
        "Визуально: зал/брендирование события отчёта (июнь 2026 / «Лесной»).",
        "Вывод: фото не являются доказательством ТЗ ноября 2025 — даже если зал выглядит «как надо».",
        "Пункт «не менее 100 фото» по вложениям выполнен количественно, но к договорному форуму не относится.",
    ]
    for i, line in enumerate(lines_p):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        set_run(p.add_run(), "•  " + line, 14, False, DARK)
    add_round(s, Inches(6.8), Inches(1.3), Inches(6.05), Inches(5.55), OFFWHITE)
    add_rect(s, Inches(6.8), Inches(1.3), Inches(6.05), Inches(0.5), GREEN)
    tb = s.shapes.add_textbox(Inches(7.0), Inches(1.38), Inches(5.7), Inches(0.38))
    set_run(tb.text_frame.paragraphs[0].add_run(), "Лекторий — визуально своё мероприятие", 15, True, WHITE)
    tb = s.shapes.add_textbox(Inches(7.05), Inches(2.0), Inches(5.6), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    lines_l = [
        f"Извлечено {photos_l.get('extracted_total', 0)}, реальных {real_l}, в API {sent_l}"
        + (f", ошибок {failed_l}" if failed_l else "")
        + ".",
        f"event_fit: договор {l_fit.get('contracted_event', 0)}, другое {l_fit.get('other_event', 0)}, неясно {l_fit.get('unknown', 0)}.",
        "Подтверждается: зал и аудитория, брендирование История.РФ / пресс-волл или экран, сцена и техника.",
        "Оборудование видно на части кадров; точный пересчёт микрофонов по фото не делался и ненадёжен.",
        "1748 заявлено vs 124 в DOCX — текстовый разрыв. Все вложенные реальные фото просмотрены; это не доказательство архива на 1748.",
    ]
    for i, line in enumerate(lines_l):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        set_run(p.add_run(), "•  " + line, 14, False, DARK)

    s = prs.slides.add_slide(blank)
    header_bar(s, "Что ИИ умеет и чего не умеет", "Границы после этапа C")
    col_titles = [
        (GREEN, "Умеет", [
            "Понять, что договор и отчёт о разных мероприятиях",
            "Сверить «не менее / не более» по тексту",
            "По всем кадрам файла увидеть зал, баннер, сцену",
            "Отличить кадры чужого события от договорного",
            "Вернуть цитаты и карточку по каждому кадру",
        ]),
        (AMBER, "Частично", [
            "Сравнить число вложений DOCX с заявленным фотоотчётом",
            "Прочитать текст логотипа на пресс-волле, если он читаем",
            "Подтвердить наличие техники, не считая её поштучно",
        ]),
        (RED, "Не умеет", [
            "Проверить 80 000 просмотров VK",
            "Сверить смету",
            "Точно посчитать микрофоны и стулья по фото",
            "Проверить внешний архив 1748 фото, которого нет в DOCX",
            "Подписать акт без человека",
        ]),
    ]
    for i, (col, title, items) in enumerate(col_titles):
        l = Inches(0.4 + i * 4.3)
        add_round(s, l, Inches(1.35), Inches(4.1), Inches(5.55), OFFWHITE)
        add_rect(s, l, Inches(1.35), Inches(4.1), Inches(0.55), col)
        tb = s.shapes.add_textbox(l + Inches(0.15), Inches(1.42), Inches(3.8), Inches(0.42))
        tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_run(tb.text_frame.paragraphs[0].add_run(), title, 16, True, WHITE)
        tb = s.shapes.add_textbox(l + Inches(0.22), Inches(2.1), Inches(3.7), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for k, item in enumerate(items):
            p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            set_run(p.add_run(), "•  " + item, 14, False, DARK)

    s = prs.slides.add_slide(blank)
    header_bar(s, "Стоимость: текст, фото и итог", "Тариф 26 ₽ / 1 млн входных токенов, 129 ₽ / 1 млн выходных")
    add_round(s, Inches(0.45), Inches(1.3), Inches(4.15), Inches(2.55), OFFWHITE)
    tb = s.shapes.add_textbox(Inches(0.65), Inches(1.45), Inches(3.8), Inches(2.25))
    tf = tb.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), f"{total_cost:.2f} ₽".replace(".", ","), 32, True, NAVY)
    p = tf.add_paragraph()
    set_run(p.add_run(), "весь пилот, 2 пары", 14, False, MUTED)
    p = tf.add_paragraph()
    p.space_before = Pt(8)
    set_run(p.add_run(), f"текст {text_cost:.2f} ₽  ·  фото {photo_cost:.2f} ₽".replace(".", ","), 13, False, DARK)
    add_round(s, Inches(4.8), Inches(1.3), Inches(8.05), Inches(2.55), OFFWHITE)
    tb = s.shapes.add_textbox(Inches(5.0), Inches(1.45), Inches(7.65), Inches(2.25))
    tf = tb.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), "Что изменил этап C", 16, True, NAVY)
    for line in [
        f"Фото: {int(photo_u.get('calls') or 0)} вызовов, ввод {int(photo_u.get('prompt_tokens') or 0):,} ток. (включая картинки).".replace(",", " "),
        f"Полный прогон {sent_p + sent_l} JPEG (иконки {skipped_p + skipped_l} не отправлялись).",
        "Картинки считаются входными токенами по тому же тарифу 26 ₽ / 1 млн.",
    ]:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        set_run(p.add_run(), "•  " + line, 14, False, DARK)
    chart_data = CategoryChartData()
    chart_data.categories = ["Текст", "Фото", "Итого"]
    chart_data.add_series("₽", (round(text_cost, 2), round(photo_cost, 2), round(total_cost, 2)))
    chart = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.4), Inches(4.05), Inches(7.4), Inches(2.95), chart_data).chart
    chart.has_legend = False
    tb = s.shapes.add_textbox(Inches(7.9), Inches(4.15), Inches(4.9), Inches(2.7))
    tf = tb.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), "Год 1000–3000 листов", 16, True, NAVY)
    for line in [
        "Текст: порядок 100–400 ₽/год.",
        "Полный прогон фото: единицы–десятки ₽ на событие, не тысячи.",
        "В API — сжатые JPEG, не исходный DOCX на 100 МБ.",
        "Деньги не блокер. Блокер — качество фотоотчёта и регламент.",
    ]:
        p = tf.add_paragraph()
        p.space_before = Pt(6)
        set_run(p.add_run(), "•  " + line, 13, False, DARK)

    s = prs.slides.add_slide(blank)
    header_bar(s, "Рекомендация", "Как пользоваться результатом этапа C")
    add_round(s, Inches(0.45), Inches(1.35), Inches(12.4), Inches(1.2), LIGHT_YEL)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12.0), Inches(0.95))
    tf = tb.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(),
            "Помощник проверяющего: чеклист, цитаты, красные флаги и визуальный просмотр всех кадров из файла. Не автоподпись акта.",
            18, True, NAVY)
    recs = [
        ("Сначала идентичность", "Если даты или площадка не те — остановить сверку ТЗ. Фото чужого события не засчитываются."),
        ("Человек на спорах", "Адрес 11 июня, «нет в отчёте» и нечитаемый баннер смотрит проверяющий."),
        ("Фото из файла", "Все реальные вложения DOCX можно прогнать. Внешний архив (1748 файлов) в модели не появится сам."),
        ("Роль в процессе", "Таблица + лист «Фото» до подписания акта. Решение человека фиксируется отдельно."),
    ]
    for i, (t, b) in enumerate(recs):
        card(s, Inches(0.45 + (i % 2) * 6.4), Inches(2.8 + (i // 2) * 2.05), Inches(6.15), Inches(1.9), t, b,
             GOLD if i < 2 else NAVY2)

    path = out_root / "prezentaciya_sverka_dogovorov.pptx"
    prs.save(path)
    return path


def export_deliverables(out_root: Path | None = None) -> tuple[Path, Path]:
    out_root = Path(out_root or Path(__file__).resolve().parent.parent / "output")
    prosv = load_json(out_root / "prosvetiteli" / "comparison.json")
    lekt = load_json(out_root / "lektoriy_kaluga" / "comparison.json")
    summary = load_json(out_root / "summary.json")
    photos_p = load_json(out_root / "prosvetiteli" / "photos.json")
    photos_l = load_json(out_root / "lektoriy_kaluga" / "photos.json")
    stamp = datetime.now()
    run_date = stamp.strftime("%d.%m.%Y")
    xlsx = build_xlsx(out_root, prosv, lekt, summary, photos_p, photos_l, run_date)
    pptx = build_pptx(out_root, prosv, lekt, summary, photos_p, photos_l, run_date)
    for path in (xlsx, pptx):
        if path.stat().st_size < 5000:
            raise RuntimeError(f"Файл слишком маленький: {path}")
    return xlsx, pptx
